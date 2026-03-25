# -*- coding: utf-8 -*-
"""
generate_pptx.py
================
Script Python utilisant la bibliothèque python-pptx pour générer automatiquement
le fichier PowerPoint "Nodi_Groupe1_Formation_IA.pptx".

Ce fichier contient 22 diapositives couvrant les 6 modules du Groupe 1
de la formation IA "Нodi и светот на ВИ" (pour enfants 9–14 ans).

Usage :
    python generate_pptx.py

Dépendances :
    pip install python-pptx==0.6.23
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Constantes de couleurs (thème light / moderne)
# ---------------------------------------------------------------------------
COLOR_COVER_BG = RGBColor(0x25, 0x63, 0xEB)      # Bleu foncé — fond couverture
COLOR_MODULE_BG = RGBColor(0xEF, 0xF6, 0xFF)     # Bleu très clair — fond modules
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # Blanc
COLOR_TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)     # Quasi-noir — texte principal
COLOR_ACCENT = RGBColor(0x25, 0x63, 0xEB)         # Bleu — accent
COLOR_SUCCESS = RGBColor(0x16, 0xA3, 0x4A)        # Vert — succès
COLOR_WARNING = RGBColor(0xF9, 0x73, 0x16)        # Orange — attention/erreur
COLOR_GOLD_BG = RGBColor(0xFC, 0xD3, 0x4D)       # Doré — fond badge final
COLOR_CONTENT_BG = RGBColor(0xFF, 0xFF, 0xFF)     # Blanc — fond slides contenu

# Dimensions de la diapositive (format 16:9 standard)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Fonctions utilitaires
# ---------------------------------------------------------------------------

def set_slide_background(slide, color: RGBColor):
    """Définit la couleur de fond d'une diapositive."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text: str, left, top, width, height,
                 font_name="Calibri", font_size=18, bold=False, italic=False,
                 color=COLOR_TEXT_DARK, align=PP_ALIGN.LEFT, word_wrap=True):
    """Ajoute une zone de texte sur la diapositive et retourne le shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title(slide, text: str, color=COLOR_TEXT_DARK, font_size=32,
              top=Inches(0.4)):
    """Ajoute un titre principal à la diapositive."""
    return add_text_box(
        slide, text,
        left=Inches(0.6), top=top,
        width=Inches(12.0), height=Inches(1.1),
        font_size=font_size, bold=True,
        color=color, align=PP_ALIGN.LEFT
    )


def add_subtitle(slide, text: str, color=COLOR_TEXT_DARK, font_size=20,
                 top=Inches(1.4)):
    """Ajoute un sous-titre à la diapositive."""
    return add_text_box(
        slide, text,
        left=Inches(0.6), top=top,
        width=Inches(12.0), height=Inches(0.7),
        font_size=font_size, bold=False,
        color=color, align=PP_ALIGN.LEFT
    )


def add_bullet_list(slide, items: list, top=Inches(2.2), left=Inches(0.8),
                    width=Inches(11.5), font_size=18,
                    color=COLOR_TEXT_DARK):
    """
    Ajoute une liste de bullet points sur la diapositive.
    Chaque item est une chaîne de texte.
    """
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_divider_line(slide, top=Inches(1.55), color=COLOR_ACCENT):
    """Ajoute une ligne horizontale décorative sous le titre."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.6), top,
        Inches(12.0), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_colored_box(slide, text: str, left, top, width, height,
                    bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
                    font_size=16, bold=False):
    """Ajoute un rectangle coloré avec du texte centré (ex : badge, durée)."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


# ---------------------------------------------------------------------------
# Slide 1 — Couverture
# ---------------------------------------------------------------------------

def create_slide_01_cover(prs):
    """Slide 1 : Couverture — fond bleu foncé, texte blanc, titre centré."""
    slide_layout = prs.slide_layouts[6]  # Layout vierge
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_COVER_BG)

    # Titre principal centré
    add_text_box(
        slide, "Нodi и светот на ВИ",
        left=Inches(1.0), top=Inches(1.8),
        width=Inches(11.3), height=Inches(1.6),
        font_size=52, bold=True,
        color=COLOR_WHITE, align=PP_ALIGN.CENTER
    )

    # Sous-titre
    add_text_box(
        slide, "Група 1 — 9 до 14 години",
        left=Inches(1.0), top=Inches(3.5),
        width=Inches(11.3), height=Inches(0.8),
        font_size=28, bold=False,
        color=COLOR_WHITE, align=PP_ALIGN.CENTER
    )

    # Mention formation
    add_text_box(
        slide, "Формација за вештачка интелигенција",
        left=Inches(1.0), top=Inches(4.4),
        width=Inches(11.3), height=Inches(0.7),
        font_size=20, bold=False, italic=True,
        color=COLOR_WHITE, align=PP_ALIGN.CENTER
    )


# ---------------------------------------------------------------------------
# Slide 2 — Vue d'ensemble des modules
# ---------------------------------------------------------------------------

def create_slide_02_overview(prs):
    """Slide 2 : Vue d'ensemble des 6 modules."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Преглед на модулите", color=COLOR_ACCENT, font_size=32)
    add_divider_line(slide)

    modules = [
        "1️⃣  Нodi се буди — Што е ВИ?",
        "2️⃣  Нodi е гладен — Податоци: гориво за ВИ",
        "3️⃣  Нodi сортира — Класификација",
        "4️⃣  Нodi полага испит — Тренирање vs Тестирање",
        "5️⃣  Нodi греши — Грешки и доверба",
        "6️⃣  Нodi е мој асистент! — ВИ во секојдневниот живот",
    ]
    add_bullet_list(slide, modules, top=Inches(1.7), font_size=20,
                    color=COLOR_TEXT_DARK)


# ---------------------------------------------------------------------------
# MODULE 1 — "Нodi се буди" (Slides 3–5)
# ---------------------------------------------------------------------------

def create_slide_03_module1_title(prs):
    """Slide 3 : Titre du Module 1."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_MODULE_BG)

    # Numéro du module
    add_colored_box(
        slide, "МОДУЛ 1",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(2.2), height=Inches(0.5),
        bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
        font_size=14, bold=True
    )

    add_title(slide, "Нodi се буди", color=COLOR_ACCENT, font_size=40,
              top=Inches(1.1))
    add_subtitle(slide, "Што е вештачка интелигенција (ВИ)?",
                 color=COLOR_TEXT_DARK, font_size=24, top=Inches(2.1))

    # Durée
    add_colored_box(
        slide, "⏱ 2 мин 30 сек",
        left=Inches(0.6), top=Inches(3.0),
        width=Inches(3.0), height=Inches(0.55),
        bg_color=COLOR_SUCCESS, text_color=COLOR_WHITE,
        font_size=15, bold=False
    )


def create_slide_04_module1_story(prs):
    """Slide 4 : Histoire du Module 1."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 1 — Приказната", color=COLOR_ACCENT, font_size=28)
    add_divider_line(slide)

    bullets = [
        "🤖  Нodi се буди и не знае НИШТО",
        "❓  Гледа фотографија — мачка или куче?",
        "💡  Добива 10 примери со ознака „мачка“",
        "⭐  Учи од примери — не од магија!",
        "📱  Исто прават: Spotify, Google Photos, GPS",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.75), font_size=20,
                    color=COLOR_TEXT_DARK)


def create_slide_05_module1_keys(prs):
    """Slide 5 : Messages clés du Module 1."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 1 — Запомни 3 работи:", color=COLOR_ACCENT,
              font_size=28)
    add_divider_line(slide)

    keys = [
        "⭐  ВИ учи од ПРИМЕРИ — не е магија!",
        "⭐  ВИ е НАСЕКАДЕ — телефон, апликации, интернет",
        "⭐  ВИ може да ГРЕШИ — и тоа е нормално!",
    ]
    add_bullet_list(slide, keys, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)

    # Activité
    add_colored_box(
        slide,
        "🎮  Активност: „ВИ или не ВИ?“ — 8 картички",
        left=Inches(0.6), top=Inches(5.5),
        width=Inches(11.5), height=Inches(0.7),
        bg_color=COLOR_WARNING, text_color=COLOR_WHITE,
        font_size=16, bold=True
    )


# ---------------------------------------------------------------------------
# MODULE 2 — "Нodi е гладен" (Slides 6–8)
# ---------------------------------------------------------------------------

def create_slide_06_module2_title(prs):
    """Slide 6 : Titre du Module 2."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_MODULE_BG)

    add_colored_box(
        slide, "МОДУЛ 2",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(2.2), height=Inches(0.5),
        bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
        font_size=14, bold=True
    )

    add_title(slide, "Нodi е гладен", color=COLOR_ACCENT, font_size=40,
              top=Inches(1.1))
    add_subtitle(slide, "Податоци: гориво за ВИ",
                 color=COLOR_TEXT_DARK, font_size=24, top=Inches(2.1))

    add_colored_box(
        slide, "⏱ 2 мин 45 сек",
        left=Inches(0.6), top=Inches(3.0),
        width=Inches(3.0), height=Inches(0.55),
        bg_color=COLOR_SUCCESS, text_color=COLOR_WHITE,
        font_size=15, bold=False
    )


def create_slide_07_module2_story(prs):
    """Slide 7 : Histoire du Module 2."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 2 — Приказната", color=COLOR_ACCENT, font_size=28)
    add_divider_line(slide)

    bullets = [
        "🍎  Нodi сака да препознава овошје",
        "❌  Лоши податоци (нејасни, погрешни) → многу грешки",
        "✅  Добри податоци (јасни, разновидни) → подобро учење",
        "🍕  Аналогија: добра храна = добро чувство; лоша храна = лошо чувство",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.75), font_size=20,
                    color=COLOR_TEXT_DARK)


def create_slide_08_module2_keys(prs):
    """Slide 8 : Messages clés du Module 2."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 2 — Запомни 3 работи:", color=COLOR_ACCENT,
              font_size=28)
    add_divider_line(slide)

    keys = [
        "⭐  Повеќе примери = подобро учење",
        "⭐  Разновидни примери = поправедна ВИ",
        "⭐  Погрешни примери = погрешна ВИ",
    ]
    add_bullet_list(slide, keys, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)

    add_colored_box(
        slide,
        "🎮  Активност: „Изгради датасет!“ — глај и пушти",
        left=Inches(0.6), top=Inches(5.5),
        width=Inches(11.5), height=Inches(0.7),
        bg_color=COLOR_WARNING, text_color=COLOR_WHITE,
        font_size=16, bold=True
    )


# ---------------------------------------------------------------------------
# MODULE 3 — "Нodi сортира" (Slides 9–11)
# ---------------------------------------------------------------------------

def create_slide_09_module3_title(prs):
    """Slide 9 : Titre du Module 3."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_MODULE_BG)

    add_colored_box(
        slide, "МОДУЛ 3",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(2.2), height=Inches(0.5),
        bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
        font_size=14, bold=True
    )

    add_title(slide, "Нodi сортира", color=COLOR_ACCENT, font_size=40,
              top=Inches(1.1))
    add_subtitle(slide, "Класификација: сортирање и препознавање",
                 color=COLOR_TEXT_DARK, font_size=24, top=Inches(2.1))

    add_colored_box(
        slide, "⏱ 3 мин 00 сек",
        left=Inches(0.6), top=Inches(3.0),
        width=Inches(3.0), height=Inches(0.55),
        bg_color=COLOR_SUCCESS, text_color=COLOR_WHITE,
        font_size=15, bold=False
    )


def create_slide_10_module3_story(prs):
    """Slide 10 : Histoire du Module 3."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 3 — Приказната", color=COLOR_ACCENT, font_size=28)
    add_divider_line(slide)

    bullets = [
        "🏭  Нodi работи во фабрика за овошје",
        "📦  Задача: сортирај јаболка / портокали / банани",
        "📏  Правила: боја + форма + големина",
        "🥝  Изненадување: пристигнува КИВИ — непознат!",
        "💡  Решение: додади нови примери → нови способности",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.75), font_size=20,
                    color=COLOR_TEXT_DARK)


def create_slide_11_module3_keys(prs):
    """Slide 11 : Messages clés du Module 3."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 3 — Запомни 3 работи:", color=COLOR_ACCENT,
              font_size=28)
    add_divider_line(slide)

    keys = [
        "⭐  Класификација = сортирање по правила од примери",
        "⭐  ВИ класифицира само она за кое е тренирана",
        "⭐  Нови примери = нови способности",
    ]
    add_bullet_list(slide, keys, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)

    add_colored_box(
        slide,
        "🎮  Активност: „Сортирај со Нodi!“ — фабрика за овошје",
        left=Inches(0.6), top=Inches(5.5),
        width=Inches(11.5), height=Inches(0.7),
        bg_color=COLOR_WARNING, text_color=COLOR_WHITE,
        font_size=16, bold=True
    )


# ---------------------------------------------------------------------------
# MODULE 4 — "Нodi полага испит" (Slides 12–14)
# ---------------------------------------------------------------------------

def create_slide_12_module4_title(prs):
    """Slide 12 : Titre du Module 4."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_MODULE_BG)

    add_colored_box(
        slide, "МОДУЛ 4",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(2.2), height=Inches(0.5),
        bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
        font_size=14, bold=True
    )

    add_title(slide, "Нodi полага испит", color=COLOR_ACCENT, font_size=40,
              top=Inches(1.1))
    add_subtitle(slide, "Тренирање vs Тестирање",
                 color=COLOR_TEXT_DARK, font_size=24, top=Inches(2.1))

    add_colored_box(
        slide, "⏱ 2 мин 45 сек",
        left=Inches(0.6), top=Inches(3.0),
        width=Inches(3.0), height=Inches(0.55),
        bg_color=COLOR_SUCCESS, text_color=COLOR_WHITE,
        font_size=15, bold=False
    )


def create_slide_13_module4_story(prs):
    """Slide 13 : Histoire du Module 4."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 4 — Приказната", color=COLOR_ACCENT, font_size=28)
    add_divider_line(slide)

    bullets = [
        "📚  Нodi учи со 20 картички (тренирање)",
        "❌  Лош начин: памти само одговори → паѓа на испит",
        "✅  Добар начин: ги разбира концептите → успева",
        "🔄  Исто важи за ВИ: тренирање ≠ тестирање",
        "⚠️  Претренираност (overfitting): добар на тренинг, лош на нови примери",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.75), font_size=19,
                    color=COLOR_TEXT_DARK)


def create_slide_14_module4_keys(prs):
    """Slide 14 : Messages clés du Module 4."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 4 — Запомни 3 работи:", color=COLOR_ACCENT,
              font_size=28)
    add_divider_line(slide)

    keys = [
        "⭐  Тренирање = учење од примери",
        "⭐  Тестирање = проверка на НОВИ примери",
        "⭐  Претренираност = ВИ која само „памти“",
    ]
    add_bullet_list(slide, keys, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)

    add_colored_box(
        slide,
        "🎮  Активност: „Тренирај и тестирај Нodi!“",
        left=Inches(0.6), top=Inches(5.5),
        width=Inches(11.5), height=Inches(0.7),
        bg_color=COLOR_WARNING, text_color=COLOR_WHITE,
        font_size=16, bold=True
    )


# ---------------------------------------------------------------------------
# MODULE 5 — "Нodi греши" (Slides 15–17)
# ---------------------------------------------------------------------------

def create_slide_15_module5_title(prs):
    """Slide 15 : Titre du Module 5."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_MODULE_BG)

    add_colored_box(
        slide, "МОДУЛ 5",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(2.2), height=Inches(0.5),
        bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
        font_size=14, bold=True
    )

    add_title(slide, "Нodi греши", color=COLOR_ACCENT, font_size=40,
              top=Inches(1.1))
    add_subtitle(slide, "Грешки, доверба и кога да верувaме на ВИ",
                 color=COLOR_TEXT_DARK, font_size=24, top=Inches(2.1))

    add_colored_box(
        slide, "⏱ 3 мин 00 сек",
        left=Inches(0.6), top=Inches(3.0),
        width=Inches(3.0), height=Inches(0.55),
        bg_color=COLOR_SUCCESS, text_color=COLOR_WHITE,
        font_size=15, bold=False
    )


def create_slide_16_module5_story(prs):
    """Slide 16 : Histoire du Module 5."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 5 — Приказната", color=COLOR_ACCENT, font_size=28)
    add_divider_line(slide)

    bullets = [
        "🔍  Нodi е детектив за спам пораки",
        "❌  Лаж позитивен: блокира важна порака од мама",
        "❌  Лаж негативен: пропушта вистински спам",
        "⚖️  Секоја ВИ греши — важно е КОГА да и верувaш",
        "🏥  Опасни грешки (медицина, безбедност) → секогаш проверувај сам!",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.75), font_size=19,
                    color=COLOR_TEXT_DARK)


def create_slide_17_module5_keys(prs):
    """Slide 17 : Messages clés du Module 5."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 5 — Запомни 3 работи:", color=COLOR_ACCENT,
              font_size=28)
    add_divider_line(slide)

    keys = [
        "⭐  Лаж позитивен = грешна аларма",
        "⭐  Лаж негативен = пропуштена аларма",
        "⭐  Секогаш проверувај кога одлуката е важна!",
    ]
    add_bullet_list(slide, keys, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)

    add_colored_box(
        slide,
        "🎮  Активност: „Биди детектив!“ — сортирај пораки",
        left=Inches(0.6), top=Inches(5.5),
        width=Inches(11.5), height=Inches(0.7),
        bg_color=COLOR_WARNING, text_color=COLOR_WHITE,
        font_size=16, bold=True
    )


# ---------------------------------------------------------------------------
# MODULE 6 — "Нodi е мој асистент!" (Slides 18–20)
# ---------------------------------------------------------------------------

def create_slide_18_module6_title(prs):
    """Slide 18 : Titre du Module 6."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_MODULE_BG)

    add_colored_box(
        slide, "МОДУЛ 6",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(2.2), height=Inches(0.5),
        bg_color=COLOR_ACCENT, text_color=COLOR_WHITE,
        font_size=14, bold=True
    )

    add_title(slide, "Нodi е мој асистент!", color=COLOR_ACCENT, font_size=38,
              top=Inches(1.1))
    add_subtitle(slide, "Како ВИ може да ти помогне во учењето",
                 color=COLOR_TEXT_DARK, font_size=24, top=Inches(2.1))

    add_colored_box(
        slide, "⏱ 3 мин 30 сек",
        left=Inches(0.6), top=Inches(3.0),
        width=Inches(3.0), height=Inches(0.55),
        bg_color=COLOR_SUCCESS, text_color=COLOR_WHITE,
        font_size=15, bold=False
    )


def create_slide_19_module6_story(prs):
    """Slide 19 : Histoire du Module 6."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 6 — Приказната", color=COLOR_ACCENT, font_size=28)
    add_divider_line(slide)

    bullets = [
        "📖  Дете има тест утре — математика (дропки)",
        "🍕  Нodi објаснува со аналогија: пица поделена на 4 дела",
        "🎯  Персонализирани прашања — прилагодени на нивото",
        "📚  4 употреби на ВИ: објаснување / повторување / организација / јазици",
        "🚫  Правила: не препишувај / секогаш проверувај / чувај ги личните податоци",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.75), font_size=19,
                    color=COLOR_TEXT_DARK)


def create_slide_20_module6_keys(prs):
    """Slide 20 : Messages clés du Module 6 + Badge."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Модул 6 — Запомни 3 работи:", color=COLOR_ACCENT,
              font_size=28)
    add_divider_line(slide)

    keys = [
        "⭐  ВИ е алатка — користи ја правилно!",
        "⭐  Секогаш прашај ЗОШТО — не само кој е одговорот",
        "⭐  Никогаш не давај лични информации на ВИ",
    ]
    add_bullet_list(slide, keys, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)

    # Badge de félicitations
    add_colored_box(
        slide,
        "🏆  ЧЕСТИТКИ! Го заврши Група 1!  |  „ИСТРАЖУВАЧ НА ВИ ⭐“",
        left=Inches(0.6), top=Inches(5.3),
        width=Inches(11.5), height=Inches(0.85),
        bg_color=COLOR_GOLD_BG, text_color=COLOR_TEXT_DARK,
        font_size=17, bold=True
    )


# ---------------------------------------------------------------------------
# Slide 21 — Évaluation finale
# ---------------------------------------------------------------------------

def create_slide_21_quiz(prs):
    """Slide 21 : Quiz final — Groupe 1."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CONTENT_BG)

    add_title(slide, "Завршен Квиз — Група 1", color=COLOR_ACCENT,
              font_size=32)
    add_divider_line(slide)

    items = [
        "📝  15 прашања од сите 6 модули",
        "✅  Потребно: минимум 70% точни одговори",
        "🏆  Награда: Значка „Истражувач на ВИ“",
        "⏱  Времетраење: околу 20 минути",
    ]
    add_bullet_list(slide, items, top=Inches(1.75), font_size=22,
                    color=COLOR_TEXT_DARK)


# ---------------------------------------------------------------------------
# Slide 22 — Badge final
# ---------------------------------------------------------------------------

def create_slide_22_badge(prs):
    """Slide 22 : Badge final — fond doré."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_GOLD_BG)

    # Grand titre badge
    add_text_box(
        slide, "🏆  ИСТРАЖУВАЧ НА ВИ",
        left=Inches(1.0), top=Inches(0.5),
        width=Inches(11.3), height=Inches(1.1),
        font_size=40, bold=True,
        color=COLOR_TEXT_DARK, align=PP_ALIGN.CENTER
    )

    # Sous-titre
    add_text_box(
        slide, "Група 1 — Завршена со успех!",
        left=Inches(1.0), top=Inches(1.55),
        width=Inches(11.3), height=Inches(0.65),
        font_size=22, bold=False,
        color=COLOR_TEXT_DARK, align=PP_ALIGN.CENTER
    )

    # Liste des modules complétés
    modules_done = [
        "✅  Модул 1: Што е ВИ?",
        "✅  Модул 2: Податоци",
        "✅  Модул 3: Класификација",
        "✅  Модул 4: Тренирање vs Тестирање",
        "✅  Модул 5: Грешки и доверба",
        "✅  Модул 6: ВИ во секојдневниот живот",
    ]
    add_bullet_list(slide, modules_done, top=Inches(2.35),
                    left=Inches(3.5), width=Inches(6.5),
                    font_size=18, color=COLOR_TEXT_DARK)


# ---------------------------------------------------------------------------
# Fonction principale
# ---------------------------------------------------------------------------

def generate_presentation():
    """Génère le fichier PPTX complet avec les 22 diapositives."""
    prs = Presentation()

    # Définir les dimensions de la diapositive (16:9)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Génération des diapositives...")

    # --- Diapositive 1 : Couverture ---
    create_slide_01_cover(prs)
    print("  [1/22] Couverture ✓")

    # --- Diapositive 2 : Vue d'ensemble ---
    create_slide_02_overview(prs)
    print("  [2/22] Vue d'ensemble des modules ✓")

    # --- Module 1 (slides 3–5) ---
    create_slide_03_module1_title(prs)
    print("  [3/22] Module 1 — Titre ✓")
    create_slide_04_module1_story(prs)
    print("  [4/22] Module 1 — Histoire ✓")
    create_slide_05_module1_keys(prs)
    print("  [5/22] Module 1 — Messages clés ✓")

    # --- Module 2 (slides 6–8) ---
    create_slide_06_module2_title(prs)
    print("  [6/22] Module 2 — Titre ✓")
    create_slide_07_module2_story(prs)
    print("  [7/22] Module 2 — Histoire ✓")
    create_slide_08_module2_keys(prs)
    print("  [8/22] Module 2 — Messages clés ✓")

    # --- Module 3 (slides 9–11) ---
    create_slide_09_module3_title(prs)
    print("  [9/22] Module 3 — Titre ✓")
    create_slide_10_module3_story(prs)
    print("  [10/22] Module 3 — Histoire ✓")
    create_slide_11_module3_keys(prs)
    print("  [11/22] Module 3 — Messages clés ✓")

    # --- Module 4 (slides 12–14) ---
    create_slide_12_module4_title(prs)
    print("  [12/22] Module 4 — Titre ✓")
    create_slide_13_module4_story(prs)
    print("  [13/22] Module 4 — Histoire ✓")
    create_slide_14_module4_keys(prs)
    print("  [14/22] Module 4 — Messages clés ✓")

    # --- Module 5 (slides 15–17) ---
    create_slide_15_module5_title(prs)
    print("  [15/22] Module 5 — Titre ✓")
    create_slide_16_module5_story(prs)
    print("  [16/22] Module 5 — Histoire ✓")
    create_slide_17_module5_keys(prs)
    print("  [17/22] Module 5 — Messages clés ✓")

    # --- Module 6 (slides 18–20) ---
    create_slide_18_module6_title(prs)
    print("  [18/22] Module 6 — Titre ✓")
    create_slide_19_module6_story(prs)
    print("  [19/22] Module 6 — Histoire ✓")
    create_slide_20_module6_keys(prs)
    print("  [20/22] Module 6 — Messages clés + Badge ✓")

    # --- Évaluation finale (slide 21) ---
    create_slide_21_quiz(prs)
    print("  [21/22] Évaluation finale ✓")

    # --- Badge final (slide 22) ---
    create_slide_22_badge(prs)
    print("  [22/22] Badge final ✓")

    # Sauvegarde du fichier
    output_file = "Nodi_Groupe1_Formation_IA.pptx"
    prs.save(output_file)
    print(f"\n✅ Fichier généré avec succès : {output_file}")
    print(f"   → 22 diapositives | 6 modules | Groupe 1")


if __name__ == "__main__":
    generate_presentation()
